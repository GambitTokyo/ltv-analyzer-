"""PPTX Export Module v245"""
import io, copy, os, math
from math import gamma
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
from matplotlib.lines import Line2D
import matplotlib.font_manager as fm
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from lxml import etree
from datetime import date
from lang import fmt_c, cur_symbol, T, get_lang, BIZ_SUBSCRIPTION, BIZ_SPOT

_JP_FP = None
def _init():
    global _JP_FP
    _here = os.path.dirname(os.path.abspath(__file__))
    candidates = [
        os.path.join(_here, 'ipag.ttf'),
        '/usr/share/fonts/opentype/ipafont-gothic/ipag.ttf',
        '/usr/share/fonts/opentype/ipafont-gothic/ipagp.ttf',
        '/usr/share/fonts/truetype/fonts-japanese-gothic.ttf',
        '/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc',
    ]
    for p in candidates:
        if os.path.exists(p):
            fm.fontManager.addfont(p); _JP_FP = fm.FontProperties(fname=p)
            plt.rcParams['font.family'] = _JP_FP.get_name()
            plt.rcParams['axes.unicode_minus'] = False; return
    for p in fm.findSystemFonts():
        if any(k in p.lower() for k in ['cjk','ipag','japanese','gothic']):
            fm.fontManager.addfont(p); _JP_FP = fm.FontProperties(fname=p)
            plt.rcParams['font.family'] = _JP_FP.get_name()
            plt.rcParams['axes.unicode_minus'] = False; return
_init()
BG = '#111820'
def _fp(): return _JP_FP or fm.FontProperties()
A = 'http://schemas.openxmlformats.org/drawingml/2006/main'

def _set_text(sh, txt):
    if not sh.has_text_frame: return
    tf = sh.text_frame
    for p in tf.paragraphs:
        for r in p.runs: r.text = ''
    if tf.paragraphs and tf.paragraphs[0].runs:
        tf.paragraphs[0].runs[0].text = txt
    elif tf.paragraphs and txt:
        # runがない場合: endParaRPrの書式を引き継いでrunを新規作成
        p0 = tf.paragraphs[0]._p
        epr = p0.find(f'{{{A}}}endParaRPr')
        import copy as _cp
        r = etree.SubElement(p0, f'{{{A}}}r')
        if epr is not None:
            rPr = _cp.deepcopy(epr)
            rPr.tag = f'{{{A}}}rPr'
            r.insert(0, rPr)
        t = etree.SubElement(r, f'{{{A}}}t')
        t.text = txt
        # runをendParaRPrの前に移動
        if epr is not None:
            p0.remove(r)
            p0.insert(list(p0).index(epr), r)

def _replace_image(sld, sh, buf):
    blip = sh._element.find('.//' + qn('a:blip'))
    if blip is None: return
    buf.seek(0); ip, rId = sld.part.get_or_add_image_part(buf); blip.set(qn('r:embed'), rId)

def _write_table(tbl, hdr, rows, footer=None):
    # ヘッダー行を上書き
    if hdr and len(tbl.rows) > 0:
        for ci, v in enumerate(hdr):
            if ci >= len(tbl.rows[0].cells): break
            tf = tbl.rows[0].cells[ci].text_frame
            for p in tf.paragraphs:
                for r in p.runs: r.text = ''
            if tf.paragraphs and tf.paragraphs[0].runs: tf.paragraphs[0].runs[0].text = str(v)
    data = rows + ([footer] if footer else [])
    need = max(2, 1 + len(data))
    tr = copy.deepcopy(tbl.rows[1]._tr) if len(tbl.rows) > 1 else None
    while len(tbl.rows) < need and tr: tbl._tbl.append(copy.deepcopy(tr))
    while len(tbl.rows) > need and len(tbl.rows) > 1: tbl._tbl.remove(tbl.rows[len(tbl.rows)-1]._tr)
    for ri, rd in enumerate(data):
        if ri+1 >= len(tbl.rows): break
        row = tbl.rows[ri+1]
        for ci, v in enumerate(rd):
            if ci >= len(row.cells): break
            tf = row.cells[ci].text_frame
            for p in tf.paragraphs:
                for r in p.runs: r.text = ''
            if tf.paragraphs and tf.paragraphs[0].runs: tf.paragraphs[0].runs[0].text = str(v)

def _write_table_styled(tbl, hdr, rows, footer=None, special_last_n=1, data_font_size=None):
    _write_table(tbl, hdr, rows, footer)
    data = rows + ([footer] if footer else [])
    n_special = special_last_n if footer is None else special_last_n
    for ri in range(len(data)):
        if ri+1 >= len(tbl.rows): break
        row = tbl.rows[ri+1]
        is_special = ri >= len(data) - n_special
        bg = '142030' if is_special else ('0D1520' if ri%2==0 else '0A1018')
        fc = 'A8DADC' if is_special else 'C8D0D8'
        for ci in range(len(row.cells)):
            cell = row.cells[ci]
            tcPr = cell._tc.get_or_add_tcPr()
            for sf in tcPr.findall(f'{{{A}}}solidFill'): tcPr.remove(sf)
            sf = etree.SubElement(tcPr, f'{{{A}}}solidFill')
            c = etree.SubElement(sf, f'{{{A}}}srgbClr'); c.set('val', bg)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    rPr = r._r.find(f'{{{A}}}rPr')
                    if rPr is None: rPr = etree.SubElement(r._r, f'{{{A}}}rPr'); r._r.insert(0, rPr)
                    for sf2 in rPr.findall(f'{{{A}}}solidFill'): rPr.remove(sf2)
                    sf2 = etree.Element(f'{{{A}}}solidFill')
                    c2 = etree.SubElement(sf2, f'{{{A}}}srgbClr'); c2.set('val', fc)
                    rPr.insert(0, sf2)
                    if data_font_size is not None:
                        rPr.set('sz', str(int(data_font_size * 100)))

def _copy_slide(prs, idx):
    src = prs.slides[idx]; new = prs.slides.add_slide(src.slide_layout)
    for el in list(new.shapes._spTree): new.shapes._spTree.remove(el)
    for el in list(src.shapes._spTree): new.shapes._spTree.append(copy.deepcopy(el))
    for rId, rel in src.part.rels.items():
        if rel.reltype == RT.IMAGE: new.part.rels._rels[rId] = rel
    # コピーされたスライドに明示的にダーク背景を設定
    P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    cSld = new._element.find(f'{{{P}}}cSld')
    if cSld is not None:
        # 既存のbgを削除
        for bg in cSld.findall(f'{{{P}}}bg'): cSld.remove(bg)
        # ダーク背景を追加（#0A0E14）
        bg = etree.SubElement(cSld, f'{{{P}}}bg')
        bgPr = etree.SubElement(bg, f'{{{P}}}bgPr')
        sf = etree.SubElement(bgPr, f'{{{A}}}solidFill')
        c = etree.SubElement(sf, f'{{{A}}}srgbClr'); c.set('val', '0A0E14')
        etree.SubElement(bgPr, f'{{{A}}}effectLst')
        # bgをcSldの最初の子要素に移動（PPTXの仕様上bgはspTreeの前にあるべき）
        cSld.remove(bg)
        cSld.insert(0, bg)
    return new

def _set_note_text(sh, note_text):
    """NOTE: テンプレートrunの書式を保持しつつテキスト差し替え+色修正"""
    if not sh.has_text_frame: return
    p0 = sh.text_frame.paragraphs[0] if sh.text_frame.paragraphs else None
    if p0 is None: return
    runs = list(p0.runs)
    if ' — ' in note_text:
        prefix, body = note_text.split(' — ', 1)
        body = '\xa0— ' + body
    else:
        prefix, body = note_text, ''
    if len(runs) >= 2:
        runs[0].text = prefix
        rPr0 = runs[0]._r.find(f'{{{A}}}rPr')
        if rPr0 is not None:
            for sf in rPr0.findall(f'{{{A}}}solidFill'): rPr0.remove(sf)
            sf = etree.Element(f'{{{A}}}solidFill')
            c = etree.SubElement(sf, f'{{{A}}}srgbClr'); c.set('val', '3A6A7A')
            rPr0.insert(0, sf)
        runs[1].text = body
        rPr1 = runs[1]._r.find(f'{{{A}}}rPr')
        if rPr1 is not None:
            rPr1.set('b', '0')
            if 'cap' in rPr1.attrib: del rPr1.attrib['cap']
            rPr1.set('lang', 'ja-JP')
            for sf in rPr1.findall(f'{{{A}}}solidFill'): rPr1.remove(sf)
            sf = etree.Element(f'{{{A}}}solidFill')
            c = etree.SubElement(sf, f'{{{A}}}srgbClr'); c.set('val', '888888')
            rPr1.insert(0, sf)
            for latin in rPr1.findall(f'{{{A}}}latin'): rPr1.remove(latin)
            for ea in rPr1.findall(f'{{{A}}}ea'): rPr1.remove(ea)
            for eff in rPr1.findall(f'{{{A}}}effectLst'): rPr1.remove(eff)
        # 3つ目以降の空runをXMLから完全削除
        p0_xml = p0._p
        all_runs = p0_xml.findall(f'{{{A}}}r')
        for r_el in all_runs[2:]:
            p0_xml.remove(r_el)
    elif len(runs) == 1:
        runs[0].text = prefix + body

def _set_s4_guide(sh, g, cur='JPY'):
    if not sh.has_text_frame: return
    txBody = sh.text_frame._txBody
    for p in txBody.findall(f'{{{A}}}p'): txBody.remove(p)
    def _para(text, color='C8D0D8', bold=False, sz=1000, indent=0):
        p = etree.SubElement(txBody, f'{{{A}}}p')
        if indent > 0:
            pPr = etree.SubElement(p, f'{{{A}}}pPr'); pPr.set('marL', str(indent)); pPr.set('indent', str(-indent))
            bc = etree.SubElement(pPr, f'{{{A}}}buClr'); c = etree.SubElement(bc, f'{{{A}}}srgbClr'); c.set('val', color)
            bs = etree.SubElement(pPr, f'{{{A}}}buSzPct'); bs.set('val', '100000')
            bu = etree.SubElement(pPr, f'{{{A}}}buChar'); bu.set('char', '•')
        r = etree.SubElement(p, f'{{{A}}}r'); rPr = etree.SubElement(r, f'{{{A}}}rPr')
        rPr.set('lang', 'ja-JP'); rPr.set('b', '1' if bold else '0'); rPr.set('sz', str(sz))
        sf = etree.SubElement(rPr, f'{{{A}}}solidFill'); c = etree.SubElement(sf, f'{{{A}}}srgbClr'); c.set('val', color)
        t = etree.SubElement(r, f'{{{A}}}t'); t.text = text
    def _mpara(segs, sz=1000, indent=0):
        p = etree.SubElement(txBody, f'{{{A}}}p')
        if indent > 0:
            pPr = etree.SubElement(p, f'{{{A}}}pPr'); pPr.set('marL', str(indent)); pPr.set('indent', str(-indent))
            bc = etree.SubElement(pPr, f'{{{A}}}buClr'); c = etree.SubElement(bc, f'{{{A}}}srgbClr'); c.set('val', segs[0][1] if segs else 'C8D0D8')
            bs = etree.SubElement(pPr, f'{{{A}}}buSzPct'); bs.set('val', '100000')
            bu = etree.SubElement(pPr, f'{{{A}}}buChar'); bu.set('char', '•')
        for text, color, bold in segs:
            r = etree.SubElement(p, f'{{{A}}}r'); rPr = etree.SubElement(r, f'{{{A}}}rPr')
            rPr.set('lang', 'ja-JP'); rPr.set('b', '1' if bold else '0'); rPr.set('sz', str(sz))
            sf = etree.SubElement(rPr, f'{{{A}}}solidFill'); c = etree.SubElement(sf, f'{{{A}}}srgbClr'); c.set('val', color)
            t = etree.SubElement(r, f'{{{A}}}t'); t.text = text
    def _empty(): etree.SubElement(txBody, f'{{{A}}}p')
    sz = 1000; ind = 228600
    _para(T('insight_title'), '3A6A7A', True, sz)
    _para(g['lam_desc'], 'C8D0D8', False, sz, ind)
    _para(g['k_desc'], 'C8D0D8', False, sz, ind)
    _mpara([(T('insight_ltv_inf', ltv=fmt_c(g['ltv_rev'], cur)), 'C8D0D8', False)], sz, ind)
    p2 = etree.SubElement(txBody, f'{{{A}}}p')
    pPr2 = etree.SubElement(p2, f'{{{A}}}pPr'); pPr2.set('marL', str(ind)); pPr2.set('indent', '0')
    etree.SubElement(pPr2, f'{{{A}}}buNone')
    for text, color, bold in [
        (T('insight_1y'), '56B4D3', False), (': ', 'C8D0D8', False),
        (f'{g["pct_1y"]:.1f}%', 'A8DADC', True), (f' ({fmt_c(g["ltv_1y"], cur)}), ', 'C8D0D8', False),
        (T('insight_2y'), '56B4D3', False), (': ', 'C8D0D8', False),
        (f'{g["pct_2y"]:.1f}%', 'A8DADC', True), (f' ({fmt_c(g["ltv_2y"], cur)}), ', 'C8D0D8', False),
        (T('insight_3y'), '56B4D3', False), (': ', 'C8D0D8', False),
        (f'{g["pct_3y"]:.1f}%', 'A8DADC', True), (f' ({fmt_c(g["ltv_3y"], cur)})', 'C8D0D8', False),
    ]:
        r = etree.SubElement(p2, f'{{{A}}}r'); rPr = etree.SubElement(r, f'{{{A}}}rPr')
        rPr.set('lang', 'ja-JP'); rPr.set('b', '1' if bold else '0'); rPr.set('sz', str(sz))
        sf = etree.SubElement(rPr, f'{{{A}}}solidFill'); c = etree.SubElement(sf, f'{{{A}}}srgbClr'); c.set('val', color)
        t = etree.SubElement(r, f'{{{A}}}t'); t.text = text
    _mpara([(f'{T("chart_cac_cap")} ({fmt_c(g["cac_upper"], cur)}): ', 'C8D0D8', False),
        (g['cac_recover_rev_str'], 'A8DADC', True), (' / GP: ', 'C8D0D8', False),
        (g['cac_recover_gp_str'], '56B4D3', True)], sz, ind)
    _empty()
    _lam_str = f'{g["lam_actual_round"]:,}{T("chart_days_suffix")}'
    _y_str = f'{g["lam_years"]:.1f}{T("chart_year_suffix")}'
    _mpara([(T('pptx_cac_design'), '56B4D3', True), ('：' if get_lang()=='ja' else ': ', 'C8D0D8', False),
        (T('pptx_cac_design_body', lam=_lam_str, y=_y_str, gp=fmt_c(g["lam_gp"], cur)), 'C8D0D8', False)], sz)

# ── グラフ（S5: 日本語に戻す） ──
def _make_ltv_graph(t_range, rev_line, gp_line, cac_line, ltv_rev, lam_actual, x_max, cur='JPY'):
    fp = _fp()
    fig, ax = plt.subplots(figsize=(10, 4), dpi=120)
    fig.patch.set_facecolor(BG); ax.set_facecolor(BG)
    marker_days = sorted(set(d for d in [180, 365, 1095, 1825, round(lam_actual)] if d <= max(t_range)))
    def _fy(line, day):
        idx = min(range(len(t_range)), key=lambda i: abs(t_range[i]-day))
        return line[idx]
    ax.plot(t_range, rev_line, color='#56b4d3', lw=1.8)
    ax.plot(t_range, gp_line, color='#a8dadc', lw=1.8, ls='--')
    ax.plot(t_range, cac_line, color='#4a7a8a', lw=1.4, ls=':')
    for line, clr in [(rev_line,'#56b4d3'),(gp_line,'#a8dadc'),(cac_line,'#4a7a8a')]:
        ax.plot(marker_days, [_fy(line,d) for d in marker_days], 'o', color=clr, ms=4, zorder=5)
    ax.axhline(ltv_rev, color='#56b4d3', lw=0.8, ls=':', alpha=0.5)
    ax.axvline(lam_actual, color='#a8dadc', lw=1.0, ls='--', alpha=0.7, ymax=0.85)
    ax.set_xlim(0, x_max+30); ax.set_ylim(0, ltv_rev*1.15)
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda v,_:f'{int(v):,}'))
    ax.set_xticks([180,365,1095,1825])
    ax.set_xticklabels([f'180{T("chart_days_suffix")}',f'1{T("chart_year_suffix")}',f'3{T("chart_year_suffix")}',f'5{T("chart_year_suffix")}'], fontproperties=fp, color='#888', fontsize=8)
    ax.tick_params(colors='#888', labelsize=8, length=0)
    for l in ax.get_yticklabels(): l.set_color('#888')
    ax.grid(True, alpha=0.15, color='#1a3040')
    for s in ax.spines.values(): s.set_color('#1a3040')
    ax.set_xlabel(T('chart_duration'), color='#888', fontsize=9, fontproperties=fp)
    ax.set_ylabel(T('chart_amount'), color='#888', fontsize=9, fontproperties=fp)
    leg = ax.legend(
        [Line2D([0],[0],color='#56b4d3',lw=1.2,marker='o',ms=2),
         Line2D([0],[0],color='#a8dadc',lw=1.2,ls='--',marker='o',ms=2),
         Line2D([0],[0],color='#4a7a8a',lw=1.0,ls=':')],
        [T('chart_ltv_rev'),T('chart_ltv_gp'),T('chart_cac_cap')],
        loc='upper left', frameon=False, fontsize=8, labelcolor='white', ncol=3,
        bbox_to_anchor=(0.02, 0.98), borderaxespad=0)
    for t in leg.get_texts(): t.set_fontproperties(fp)
    ax.annotate(T('chart_lam_days', n=int(lam_actual)), xy=(lam_actual, ltv_rev*0.92), color='#a8dadc', fontsize=8, fontproperties=fp, ha='center', annotation_clip=False, bbox=dict(boxstyle='square,pad=0.15', facecolor=BG, edgecolor='none'))
    ax.annotate(f'LTV∞ {fmt_c(ltv_rev, cur)}', xy=(x_max+32, ltv_rev), color='#56b4d3', fontsize=8, fontproperties=fp, va='center', annotation_clip=False, clip_on=False)
    fig.subplots_adjust(left=0.08, right=0.88, top=0.95, bottom=0.14)
    buf = io.BytesIO(); fig.savefig(buf, format='png', dpi=120, facecolor=BG); buf.seek(0); plt.close(fig)
    return buf

# ── S7/S13 棒グラフ：数値ラベル+X軸数値+加重平均ライン数値 ──
def _make_bar_graph(pp_rows, best, avg_ltv, cur='JPY'):
    fp = _fp()
    fig, ax = plt.subplots(figsize=(6, max(2.5, len(pp_rows)*0.45)), dpi=120)
    fig.patch.set_facecolor(BG); ax.set_facecolor(BG)
    segs = [r['seg'] for r in pp_rows][::-1]
    vals = [r['ltv_r'] for r in pp_rows][::-1]
    colors = ['#56b4d3' if pp_rows[::-1][i]['seg']==best['seg'] else '#2a4a5a' for i in range(len(pp_rows))]
    # 加重平均ラインを先に描画（背面）
    ax.axvline(avg_ltv, color='#a8dadc', lw=1, ls='--', alpha=0.7, zorder=1)
    ax.text(avg_ltv, len(segs)-0.3, f'Avg {fmt_c(avg_ltv, cur)}', color='#a8dadc', fontsize=7, ha='center', va='bottom', fontproperties=fp, zorder=5)
    # 棒グラフを後に描画（最前面）
    bars = ax.barh(segs, vals, color=colors, height=0.6, zorder=3)
    # 棒内に数値ラベル（右寄せ）
    for bar, val in zip(bars, vals):
        ax.text(bar.get_width() - max(vals)*0.02, bar.get_y() + bar.get_height()/2,
                f'{fmt_c(val, cur)}', va='center', ha='right', color='white', fontsize=7, fontproperties=fp, zorder=4)
    ax.tick_params(colors='#888', labelsize=8)
    for l in ax.get_yticklabels(): l.set_fontproperties(fp)
    # X軸: 数値表記（60,000形式）
    ax.xaxis.set_major_formatter(mticker.FuncFormatter(lambda v,_:f'{int(v):,}'))
    for l in ax.get_xticklabels(): l.set_fontproperties(fp); l.set_fontsize(7)
    for s in ax.spines.values(): s.set_color('#1a3040')
    ax.grid(True, alpha=0.15, color='#1a3040', axis='x')
    fig.tight_layout()
    buf = io.BytesIO(); fig.savefig(buf, format='png', dpi=120, facecolor=BG); buf.seek(0); plt.close(fig)
    return buf

# ════════════════════════════════════════════════════════════════
def generate_pptx(
    tmpl_path, k, lam, lam_actual, r2, arpu_daily, gpm, gp_daily, cac_n, cac_upper,
    ltv_rev, lam_rev, lam_gp, rev_99, gp_99, days_99,
    t_range, rev_line, gp_line, cac_line, x_max,
    tbl_rows, horizons, ltv_offset_days,
    ltv_horizon_offset, ltv_horizon_spot, business_type, dormancy_days,
    buf1, buf2, df, client_name, analyst_name, billing_cycle_display,
    k_summary, r2_summary, cac_label, cac_recover_rev_str, cac_recover_gp_str,
    segment_cols_input, _compute_km_df, _fit_weibull_df, ltv_inf, fmt_horizon,
    s4_guide_data=None,
    report_title=None,
    arpu_0_dorm=None,
    arpu_long=None,
    outlier_label=T('excel_none'),
    all_seg_results=None,
    cur='JPY',
):
    if k<0.7: ki=T("insight_k_early_strong", k=k, acq=T("common_acquisition"))
    elif k<1.0: ki=T("insight_k_early_mild", k=k, acq=T("common_acquisition"))
    elif k<1.5: ki=T("insight_k_late_mild", k=k, acq=T("common_acquisition"))
    else: ki=T("insight_k_late_strong", k=k, acq=T("common_acquisition"))
    if r2>=0.95: rc=T("insight_r2_high", r2=r2)
    elif r2>=0.85: rc=T("insight_r2_mid", r2=r2)
    else: rc=T("insight_r2_low", r2=r2)
    prs = Presentation(tmpl_path)
    pd_mod = __import__('pandas')
    _dc = [df['start_date']]
    if 'end_date' in df.columns and df['end_date'].notna().any(): _dc.append(df['end_date'])
    _ds = pd_mod.concat(_dc).dropna().min().strftime('%Y/%m/%d')
    _de = pd_mod.concat(_dc).dropna().max().strftime('%Y/%m/%d')
    for sh in prs.slides[0].shapes:
        if sh.name=='TextBox 4': _set_text(sh, report_title or 'Kaplan–Meier × Weibull Model')
        elif sh.name=='TextBox 5': _set_text(sh, client_name or '')
        elif sh.name=='TextBox 6': _set_text(sh, date.today().strftime(T('pdf_date_format')))
        elif sh.name=='TextBox 7': _set_text(sh, analyst_name or '')
    s2=prs.slides[1]
    for sh in s2.shapes:
        if sh.name=='タイトル 5': _set_text(sh, T('main_summary_title'))
        if sh.name=='テキスト プレースホルダー 6' and sh.has_text_frame:
            tf=sh.text_frame
            # 1行目：データ期間 | 顧客数 | 解約済み | 継続中 | 異常値の処理
            i1=f"{T("pdf_data_period")}: {_ds} – {_de} | {T("excel_customers")}: {len(df):,} | Churned: {df['event'].sum():,} | Active: {(df['event']==0).sum():,} | Outlier: {outlier_label}"
            # 2行目：ビジネスタイプ | 請求サイクル/休眠判定 | 日割り(サブスクのみ) | Daily ARPU | GPM
            if business_type == BIZ_SPOT:
                _biz_disp = T('biz_spot')
                _dorm_disp = T('common_days_unit', n=dormancy_days) if dormancy_days else T('common_days_unit', n=180)
                i2=f"{_biz_disp} | {T("excel_dormancy")}: {_dorm_disp} | Daily ARPU: {fmt_c(arpu_daily, cur, 2)} | GPM: {gpm:.0%}"
            else:
                _biz_disp = T('biz_subscription')
                i2=f"{_biz_disp} | {billing_cycle_display} | {T('excel_prorate')}: {'ON' if ltv_offset_days==0 else 'OFF'} | Daily ARPU: {fmt_c(arpu_daily, cur, 2)} | GPM: {gpm:.0%}"
            if len(tf.paragraphs)>=1:
                for r in tf.paragraphs[0].runs: r.text=''
                if tf.paragraphs[0].runs: tf.paragraphs[0].runs[0].text=i1
            if len(tf.paragraphs)>=2:
                for r in tf.paragraphs[1].runs: r.text=''
                if tf.paragraphs[1].runs: tf.paragraphs[1].runs[0].text=i2
        elif sh.name=='グループ化 26':
            kpi={'TextBox 6':f'{fmt_c(ltv_rev, cur)}','TextBox 10':f'{fmt_c(cac_upper, cur)}','TextBox 11':f'{T("chart_cac_cap")} ({cac_label})','TextBox 14':f'{k:.3f}','TextBox 18':f'{lam_actual:.0f}{T("chart_days_suffix")}','TextBox 22':f'{r2:.3f}',
                'TextBox 7':'LTV∞','TextBox 8':T('summary_rev_basis'),
                'TextBox 12':T('summary_cac_gp_basis'),
                'TextBox 16':T('summary_k_early') if k<1 else T('summary_k_late'),
                'TextBox 20':T('summary_k_desc_long'),
                'TextBox 24':T('summary_r2_note')}
            for g in sh.shapes:
                if g.name in kpi: _set_text(g, kpi[g.name])
        elif sh.name=='テキスト ボックス 49' and sh.has_text_frame:
            tf=sh.text_frame
            if len(tf.paragraphs)>=2:
                for r in tf.paragraphs[0].runs: r.text=''
                if tf.paragraphs[0].runs: tf.paragraphs[0].runs[0].text=T('summary_conclusion')
                for r in tf.paragraphs[1].runs: r.text=''
                if tf.paragraphs[1].runs: tf.paragraphs[1].runs[0].text=k_summary+r2_summary
    s3=prs.slides[2]; ld=lam+ltv_offset_days if business_type==BIZ_SPOT else lam
    for sh in s3.shapes:
        if sh.name=='タイトル 31': _set_text(sh, T('chart_reliability_title'))
        elif sh.name=='テキスト プレースホルダー 44' and sh.has_text_frame:
            for r in sh.text_frame.paragraphs[0].runs: r.text=''
            if sh.text_frame.paragraphs[0].runs: sh.text_frame.paragraphs[0].runs[0].text=T('pptx_weibull_subtitle')
    for sh in s3.shapes:
        if sh.name=='Picture 3': buf1.seek(0); _replace_image(s3,sh,buf1)
        elif sh.name=='Picture 4': buf2.seek(0); _replace_image(s3,sh,buf2)
        elif sh.name=='TextBox 7': _set_text(sh, f"k (shape) = {k:.3f}\n→ {ki}")
        elif sh.name=='TextBox 8': _set_text(sh, f"λ (scale) = {ld:.1f}{T('chart_days_suffix')} (~{ld/365:.1f}{T('chart_year_suffix')})\n→ {T('summary_k_desc_long')}\n→ {rc}")
    s4=prs.slides[3]
    for sh in s4.shapes:
        if sh.name=='タイトル 6': _set_text(sh, T('chart_interim_ltv_title')); rows_s4=[]
    for h in horizons:
        if business_type==BIZ_SPOT:
            dp=dormancy_days or 180
            _a0 = arpu_0_dorm if arpu_0_dorm is not None else arpu_daily
            _al = arpu_long if arpu_long is not None else arpu_daily
            lr=ltv_horizon_spot(k,lam,_a0,_al,h,dp); lg=ltv_horizon_spot(k,lam,_a0*gpm,_al*gpm,h,dp)
        else:
            lr=ltv_horizon_offset(k,lam,arpu_daily,h,ltv_offset_days); lg=ltv_horizon_offset(k,lam,gp_daily,h,ltv_offset_days)
        label=f'{h}{T("chart_days_suffix")}' if h<365 else T('chart_year_days', y=h//365, d=f'{h:,}')
        rows_s4.append([label,f'{fmt_c(lr, cur)}',f'{fmt_c(lg, cur)}',f'{fmt_c(lg/cac_n, cur)}',f'{lr/ltv_rev*100:.1f}%'])
    rows_s4.append([T('tbl_lam_row', n=f'{round(lam_actual):,}'),f'{fmt_c(lam_rev, cur)}',f'{fmt_c(lam_gp, cur)}',f'{fmt_c(lam_gp/cac_n, cur)}',f'{lam_rev/ltv_rev*100:.1f}%'])
    rows_s4.append([T('tbl_99pct_row', n=f'{int(days_99):,}'),f'{fmt_c(rev_99, cur)}',f'{fmt_c(gp_99, cur)}',f'{fmt_c(gp_99/cac_n, cur)}','99.0%'])
    rows_s4.append([f'LTV∞',f'{fmt_c(ltv_rev, cur)}',f'{fmt_c(cac_upper*cac_n, cur)}',f'{fmt_c(cac_upper, cur)}','100%'])
    _tbl_shape = None
    for sh in s4.shapes:
        if sh.shape_type==19:
            _tbl_shape = sh
            _write_table_styled(sh.table,[T('tbl_horizon'),T('tbl_ltv_rev'),T('tbl_ltv_gp'),T('tbl_cac_cap'),T('tbl_pct_ltv')],rows_s4,special_last_n=3,data_font_size=8)
    for sh in s4.shapes:
        if sh.name=='テキスト ボックス 3' and sh.has_text_frame:
            if _tbl_shape:
                sh.top = _tbl_shape.top + _tbl_shape.height + 600000  # テーブル底 + 隙間
            if s4_guide_data: _set_s4_guide(sh, s4_guide_data, cur)
    buf_s5=_make_ltv_graph(t_range,rev_line,gp_line,cac_line,ltv_rev,lam_actual,x_max,cur)
    for sh in prs.slides[4].shapes:
        if sh.name=='タイトル 5': _set_text(sh, T('chart_interim_ltv_title'))
        elif sh.name=='コンテンツ プレースホルダー 7':
            sh.left=int(0.5*914400); sh.top=int(1.6*914400); sh.width=int(12.2*914400); sh.height=int(4.3*914400)
            _replace_image(prs.slides[4],sh,buf_s5)
    if not segment_cols_input.strip():
        sldIdLst=prs.slides._sldIdLst
        for idx in [9,8,7,6,5]:
            if idx<len(sldIdLst):
                e=sldIdLst[idx]; rId=e.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                if rId:
                    try: prs.part.drop_rel(rId)
                    except: pass
                sldIdLst.remove(e)
    else:
        seg_cols=[c.strip() for c in segment_cols_input.split(',') if c.strip() and c.strip() in df.columns]
        for sh in prs.slides[5].shapes:
            if sh.name=='TextBox 3': _set_text(sh, T('section_segment'))
            elif sh.name=='TextBox 4': _set_text(sh,'  |  '.join(seg_cols))
        for sc in seg_cols:
            # all_seg_resultsから参照（独自Weibullフィットしない）
            pp_rows=[]
            if all_seg_results and sc in all_seg_results:
                for _, _r in all_seg_results[sc].iterrows():
                    pp_rows.append({'seg':str(_r['segment']),'n':int(_r['n_customers']),'ltv_r':_r['ltv_rev'],'ltv_g':_r['ltv_gp'],'cac':_r['cac_cap'],'k':_r['k'],'lam':_r['λ_raw'],'r2':_r['R²'],'arpu_s':_r['arpu_s'],'arpu_long_s':_r['arpu_long_s'],'arpu_0_dorm_s':_r['arpu_0_dorm_s']})
            else:
                for sv in sorted(df[sc].dropna().unique()):
                    df_s=df[df[sc]==sv]
                    if len(df_s)<10 or df_s['event'].sum()<5: continue
                    try:
                        km_s=_compute_km_df(df_s); k_s,lam_s,r2_s,_=_fit_weibull_df(km_s)
                        if k_s is None: continue
                        arpu_s=df_s['arpu_daily'].mean(); gp_s=arpu_s*gpm
                        _si=lam_s*gamma(1+1/k_s)
                        ltv_r=((_si+ltv_offset_days)*arpu_s); ltv_g=((_si+ltv_offset_days)*gp_s)
                        pp_rows.append({'seg':str(sv),'n':len(df_s),'ltv_r':ltv_r,'ltv_g':ltv_g,'cac':ltv_g/cac_n,'k':k_s,'lam':lam_s,'r2':r2_s,'arpu_s':arpu_s,'arpu_long_s':arpu_s,'arpu_0_dorm_s':arpu_s})
                    except: continue
            if not pp_rows: continue
            pp_rows.sort(key=lambda x:x['ltv_r'],reverse=True)
            best=pp_rows[0]; n_total=sum(r['n'] for r in pp_rows)
            avg_ltv=sum(r['ltv_r']*r['n'] for r in pp_rows)/n_total
            w_ltv_g=sum(r['ltv_g']*r['n'] for r in pp_rows)/n_total
            w_cac=sum(r['cac']*r['n'] for r in pp_rows)/n_total
            avg_cac=sum(r['cac']*r['n'] for r in pp_rows)/n_total
            premium=(best['ltv_r']-avg_ltv)/avg_ltv*100
            cac_diff=best['cac']-avg_cac
            cac_diff_str=f"+{fmt_c(cac_diff, cur)} higher" if cac_diff>=0 else f"{fmt_c(abs(cac_diff), cur)} lower"
            s7=_copy_slide(prs,6)
            for sh in s7.shapes:
                if sh.name=='タイトル 4': _set_text(sh,f'{sc}: LTV∞')
                elif sh.name=='テキスト プレースホルダー 5' and sh.has_text_frame:
                    tf=sh.text_frame
                    if len(tf.paragraphs)>=2:
                        for r in tf.paragraphs[0].runs: r.text=''
                        if tf.paragraphs[0].runs: tf.paragraphs[0].runs[0].text=f'TOP PICK  {best["seg"]}'
                        for r in tf.paragraphs[1].runs: r.text=''
                        if tf.paragraphs[1].runs: tf.paragraphs[1].runs[0].text=f'LTV∞ (Rev): {fmt_c(best["ltv_r"], cur)} (vs avg +{premium:.1f}%) | {T("chart_cac_cap")} {fmt_c(best["cac"], cur)} ({cac_diff_str})'
                elif sh.name=='コンテンツ プレースホルダー 8':
                    buf7=_make_bar_graph(pp_rows[:10],best,avg_ltv,cur); _replace_image(s7,sh,buf7)
            s8=_copy_slide(prs,7); top10=pp_rows[:10]
            dr8=[[r['seg'],f'{r["n"]:,}',f'{fmt_c(r["ltv_r"], cur)}',f'{fmt_c(r["ltv_g"], cur)}',f'{fmt_c(r["cac"], cur)}',f'{r["k"]:.3f}',f'{r["lam"]:.1f}',f'{r["r2"]:.3f}'] for r in top10]
            ft8=[T('seg_weighted_avg'),f'{n_total:,}',f'{fmt_c(avg_ltv, cur)}',f'{fmt_c(w_ltv_g, cur)}',f'{fmt_c(w_cac, cur)}','—','—','—']
            _seg_hdr8=[T('seg_tbl_segment'),T('seg_tbl_n'),T('seg_tbl_ltv_rev'),T('seg_tbl_ltv_gp'),T('seg_tbl_cac_cap'),'k',T('seg_tbl_lam'),'R²']
            for sh in s8.shapes:
                if sh.name=='タイトル 2': _set_text(sh,f'{sc}: {T("pptx_summary")}')
                elif sh.shape_type==19: _write_table_styled(sh.table,_seg_hdr8,dr8,ft8)
                elif sh.name=='テキスト ボックス 9' and sh.has_text_frame:
                    diff_pct=(avg_ltv-ltv_rev)/ltv_rev*100
                    note=T('pdf_note_summary_table', max=10, total=len(pp_rows)) + f' (Δ{diff_pct:+.1f}%)'
                    _set_note_text(sh,note)
            for ri,row in enumerate(pp_rows):
                tmpl_idx=8 if ri==0 else 9
                sx=_copy_slide(prs,tmpl_idx)
                buf_km=buf_wb=None
                _arpu_s=row['arpu_s']; _arpu_long_s=row['arpu_long_s']; _arpu_0_dorm_s=row['arpu_0_dorm_s']
                _dorm_s=dormancy_days or 180 if business_type==BIZ_SPOT else ltv_offset_days
                try:
                    df_s=df[df[sc]==row['seg']]
                    # グラフ描画用KM（数値には使わない）
                    if ltv_offset_days > 0:
                        _df_s_fit=df_s.copy(); _df_s_fit['duration']=_df_s_fit['duration']-ltv_offset_days
                        _df_s_fit.loc[_df_s_fit['duration']<=0,'event']=0; _df_s_fit.loc[_df_s_fit['duration']<=0,'duration']=1
                        _df_s_fit['duration']=_df_s_fit['duration'].clip(lower=1)
                        km_s=_compute_km_df(_df_s_fit)
                    else:
                        km_s=_compute_km_df(df_s)
                    fig_km,ax_km=plt.subplots(figsize=(4,2.8),dpi=120); fig_km.patch.set_facecolor(BG); ax_km.set_facecolor(BG)
                    ax_km.step(km_s['time'],km_s['survival'],where='post',color='#56b4d3',lw=1.5)
                    t_fit=list(range(int(km_s['time'].max()))); s_fit=[math.exp(-(t/row['lam'])**row['k']) for t in t_fit]
                    ax_km.plot(t_fit,s_fit,color='#a8dadc',lw=1.2,ls='--')
                    ax_km.tick_params(colors='#888',labelsize=7)
                    for s in ax_km.spines.values(): s.set_color('#1a3040')
                    ax_km.grid(True,alpha=0.15,color='#1a3040'); fig_km.tight_layout()
                    buf_km=io.BytesIO(); fig_km.savefig(buf_km,format='png',dpi=120,facecolor=BG); buf_km.seek(0); plt.close(fig_km)
                    x_s=list(range(1,int(max(1825,row['lam']*2)),max(1,int(row['lam']*2)//200)))
                    if business_type==BIZ_SPOT:
                        y_s=[ltv_horizon_spot(row['k'],row['lam'],_arpu_0_dorm_s,_arpu_long_s,t,_dorm_s) for t in x_s]
                    else:
                        y_s=[ltv_horizon_offset(row['k'],row['lam'],_arpu_s,t,ltv_offset_days) for t in x_s]
                    fig_wb,ax_wb=plt.subplots(figsize=(4,2.8),dpi=120); fig_wb.patch.set_facecolor(BG); ax_wb.set_facecolor(BG)
                    ax_wb.plot(x_s,y_s,color='#56b4d3',lw=1.5)
                    ax_wb.axvline(row['lam'],color='#a8dadc',lw=1,ls='--',alpha=0.7)
                    ax_wb.axhline(row['ltv_r'],color='#56b4d3',lw=0.8,ls=':',alpha=0.5)
                    ax_wb.tick_params(colors='#888',labelsize=7)
                    for s in ax_wb.spines.values(): s.set_color('#1a3040')
                    ax_wb.grid(True,alpha=0.15,color='#1a3040'); fig_wb.tight_layout()
                    buf_wb=io.BytesIO(); fig_wb.savefig(buf_wb,format='png',dpi=120,facecolor=BG); buf_wb.seek(0); plt.close(fig_wb)
                except: pass
                rows_sx=[]
                for h in horizons:
                    if business_type==BIZ_SPOT:
                        lr_s=ltv_horizon_spot(row['k'],row['lam'],_arpu_0_dorm_s,_arpu_long_s,h,_dorm_s)
                        lg_s=ltv_horizon_spot(row['k'],row['lam'],_arpu_0_dorm_s*gpm,_arpu_long_s*gpm,h,_dorm_s)
                    else:
                        lr_s=ltv_horizon_offset(row['k'],row['lam'],_arpu_s,h,ltv_offset_days)
                        lg_s=lr_s*gpm
                    rows_sx.append([fmt_horizon(h),f'{fmt_c(lr_s, cur)}',f'{fmt_c(lg_s, cur)}',f'{fmt_c(lg_s/cac_n, cur)}',f'{lr_s/row["ltv_r"]*100:.1f}%'])
                # λ行
                _lam_actual_s=row['lam']+ltv_offset_days
                if business_type==BIZ_SPOT:
                    lam_r_s=ltv_horizon_spot(row['k'],row['lam'],_arpu_0_dorm_s,_arpu_long_s,_lam_actual_s,_dorm_s)
                    lam_g_s=ltv_horizon_spot(row['k'],row['lam'],_arpu_0_dorm_s*gpm,_arpu_long_s*gpm,_lam_actual_s,_dorm_s)
                else:
                    lam_r_s=ltv_horizon_offset(row['k'],row['lam'],_arpu_s,_lam_actual_s,ltv_offset_days)
                    lam_g_s=lam_r_s*gpm
                rows_sx.append([T('tbl_lam_row', n=f'{round(row["lam"]):,}'),f'{fmt_c(lam_r_s, cur)}',f'{fmt_c(lam_g_s, cur)}',f'{fmt_c(lam_g_s/cac_n, cur)}',f'{lam_r_s/row["ltv_r"]*100:.1f}%'])
                try:
                    from scipy.optimize import brentq as _bq
                    if business_type==BIZ_SPOT:
                        d99s=_bq(lambda h:ltv_horizon_spot(row['k'],row['lam'],_arpu_0_dorm_s,_arpu_long_s,h,_dorm_s)/row['ltv_r']-0.99,1,500000)
                        r99s=ltv_horizon_spot(row['k'],row['lam'],_arpu_0_dorm_s,_arpu_long_s,d99s,_dorm_s)
                        g99s=ltv_horizon_spot(row['k'],row['lam'],_arpu_0_dorm_s*gpm,_arpu_long_s*gpm,d99s,_dorm_s)
                    else:
                        d99s=_bq(lambda h:ltv_horizon_offset(row['k'],row['lam'],_arpu_s,h,ltv_offset_days)/row['ltv_r']-0.99,1,500000)
                        r99s=ltv_horizon_offset(row['k'],row['lam'],_arpu_s,d99s,ltv_offset_days)
                        g99s=r99s*gpm
                    rows_sx.append([T('tbl_99pct_row', n=f'{int(d99s):,}'),f'{fmt_c(r99s, cur)}',f'{fmt_c(g99s, cur)}',f'{fmt_c(g99s/cac_n, cur)}','99.0%'])
                except:
                    r99_approx = row['ltv_r'] * 0.99
                    g99_approx = row['ltv_g']* 0.99
                    rows_sx.append([T('tbl_99pct_row', n='~'),f'{fmt_c(r99_approx, cur)}',f'{fmt_c(g99_approx, cur)}',f'{fmt_c(g99_approx/cac_n, cur)}','99.0%'])
                # LTV∞行
                rows_sx.append([f'LTV∞',f'{fmt_c(row["ltv_r"], cur)}',f'{fmt_c(row["ltv_g"], cur)}',f'{fmt_c(row["ltv_g"]/cac_n, cur)}','100%'])
                for sh in sx.shapes:
                    if sh.name=='タイトル 8': _set_text(sh,f'{sc}: {row["seg"]}')
                    elif sh.name=='テキスト プレースホルダー 9' and sh.has_text_frame:
                        for _r9 in sh.text_frame.paragraphs[0].runs: _r9.text=''
                        if sh.text_frame.paragraphs[0].runs: sh.text_frame.paragraphs[0].runs[0].text=f'{T("chart_reliability_title")}  |  {T("chart_interim_ltv_title")}'
                    elif sh.name=='Picture 6' and buf_km: buf_km.seek(0); _replace_image(sx,sh,buf_km)
                    elif sh.name=='Picture 7' and buf_wb: buf_wb.seek(0); _replace_image(sx,sh,buf_wb)
                    elif sh.shape_type==19: _write_table_styled(sh.table,[T('tbl_horizon'),T('tbl_ltv_rev'),T('tbl_ltv_gp'),T('tbl_cac_cap'),T('tbl_pct_ltv')],rows_sx,special_last_n=3,data_font_size=8)
                    elif sh.name=='テキスト ボックス 17' and sh.has_text_frame:
                        _set_text(sh,'TOP PICK  ' if ri==0 else '')
                        bp=sh.text_frame._txBody.find(f'{{{A}}}bodyPr')
                        if bp is None: bp=etree.SubElement(sh.text_frame._txBody,f'{{{A}}}bodyPr')
                        bp.set('anchor','ctr')
        sldIdLst=prs.slides._sldIdLst
        for idx in [9,8,7,6]:
            if idx<len(sldIdLst):
                e=sldIdLst[idx]; rId=e.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                if rId:
                    try: prs.part.drop_rel(rId)
                    except: pass
                sldIdLst.remove(e)
    buf_out=io.BytesIO(); prs.save(buf_out); buf_out.seek(0)
    return buf_out
