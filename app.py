import streamlit as st
import pandas as pd
import numpy as np
from scipy import stats
from scipy.special import gamma, gammainc
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib import rcParams
import io
import warnings
warnings.filterwarnings('ignore')

# ── Page config ──────────────────────────────────────────────
st.set_page_config(
    page_title="LTV Analyzer",
    page_icon="📈",
    layout="wide",
    initial_sidebar_state="expanded"
)

# ── CSS ───────────────────────────────────────────────────────
st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Syne:wght@400;700&family=DM+Sans:wght@300;400;500&display=swap');
html, body, [class*="css"] { font-family: 'DM Sans', sans-serif; }
.stApp { background-color: #0d0d0d; color: #e8e4dc; }
.metric-card {
    background: #161616;
    border: 1px solid #242424;
    border-radius: 12px;
    padding: 20px 16px;
    text-align: center;
    height: 100%;
}
.metric-value {
    font-family: 'Syne', sans-serif;
    font-size: 2rem;
    color: #c8b89a;
    line-height: 1.1;
    word-break: break-all;
}
.metric-label {
    font-size: 0.7rem;
    color: #555;
    text-transform: uppercase;
    letter-spacing: 0.12em;
    margin-top: 6px;
}
.section-title {
    font-family: 'Syne', sans-serif;
    font-size: 1.1rem;
    color: #e8e4dc;
    border-bottom: 1px solid #222;
    padding-bottom: 6px;
    margin: 28px 0 16px 0;
}
.prompt-box {
    background: #161616;
    border: 1px solid #2a2a2a;
    border-left: 3px solid #c8b89a;
    border-radius: 8px;
    padding: 16px;
    font-family: 'DM Sans', monospace;
    font-size: 0.85rem;
    color: #ccc;
    white-space: pre-wrap;
    word-break: break-word;
}
.help-box {
    background: #141414;
    border: 1px solid #1e1e1e;
    border-radius: 8px;
    padding: 14px 16px;
    font-size: 0.82rem;
    color: #777;
    margin-top: 8px;
}
.tag {
    display: inline-block;
    background: #1e1e1e;
    border: 1px solid #2a2a2a;
    border-radius: 20px;
    padding: 2px 10px;
    font-size: 0.7rem;
    color: #888;
    margin: 2px;
}
</style>
""", unsafe_allow_html=True)

# ── Matplotlib theme ──────────────────────────────────────────
plt.style.use('dark_background')
import matplotlib.font_manager as _fm, subprocess as _sp
try:
    _jp = [f.name for f in _fm.fontManager.ttflist
           if any(x in f.name for x in ['Noto Sans CJK', 'IPAexGothic', 'Hiragino'])]
    if not _jp:
        _sp.run(['apt-get','install','-y','-q','fonts-noto-cjk'], capture_output=True)
        _fm.fontManager.__init__()
        _jp = [f.name for f in _fm.fontManager.ttflist if 'Noto Sans CJK' in f.name]
    rcParams['font.family'] = _jp[0] if _jp else 'DejaVu Sans'
except Exception:
    rcParams['font.family'] = 'DejaVu Sans'
for _k, _v in {
    'figure.facecolor': '#161616', 'axes.facecolor': '#161616',
    'axes.edgecolor': '#222', 'axes.labelcolor': '#888',
    'xtick.color': '#555', 'ytick.color': '#555',
    'grid.color': '#222', 'grid.linewidth': 0.5,
    'axes.unicode_minus': False,
}.items():
    rcParams[_k] = _v

ACCENT  = '#c8b89a'
ACCENT2 = '#7a9e9f'
ACCENT3 = '#a0856c'

# ══════════════════════════════════════════════════════════════
# Analysis functions
# ══════════════════════════════════════════════════════════════

def compute_km(df):
    df = df.sort_values('duration').reset_index(drop=True)
    S, records = 1.0, []
    for t in np.sort(df['duration'].unique()):
        n_risk   = (df['duration'] >= t).sum()
        n_events = ((df['duration'] == t) & (df['event'] == 1)).sum()
        if n_risk > 0:
            S *= (1 - n_events / n_risk)
        records.append({'t': t, 'S': S, 'n_events': n_events, 'n_risk': n_risk})
    return pd.DataFrame(records)

def fit_weibull(km_df):
    fd = km_df[(km_df['t'] > 0) & (km_df['S'] > 0) & (km_df['S'] < 1)].copy()
    fd['ln_t']         = np.log(fd['t'])
    fd['ln_neg_ln_S']  = np.log(-np.log(fd['S']))
    fd = fd.replace([np.inf, -np.inf], np.nan).dropna()
    if len(fd) < 3:
        return None, None, None, None
    k, b, r, *_ = stats.linregress(fd['ln_t'], fd['ln_neg_ln_S'])
    lam = np.exp(-b / k)
    return k, lam, r**2, fd

def ltv_inf(k, lam, arpu):
    surv_int = lam * gamma(1 + 1/k)
    return surv_int * arpu, surv_int

def ltv_horizon(k, lam, arpu, h):
    x = (h / lam) ** k
    return lam * gamma(1 + 1/k) * gammainc(1 + 1/k, x) * arpu

def weibull_s(t, k, lam):
    return np.exp(-(t / lam) ** k)

# ══════════════════════════════════════════════════════════════
# Sidebar
# ══════════════════════════════════════════════════════════════

with st.sidebar:
    st.markdown("### 📥 データ入力")

    # 1万件のリアルなサンプルデータ生成
    np.random.seed(42)
    n_sample = 10000
    start_dates = pd.date_range('2022-01-01', '2024-06-30', periods=n_sample)
    survival_days = np.random.weibull(1.2, n_sample) * 400
    churned = np.random.random(n_sample) < 0.65
    today_ts = pd.Timestamp.today()
    end_dates = []
    for i in range(n_sample):
        if churned[i]:
            ed = start_dates[i] + pd.Timedelta(days=int(survival_days[i]))
            end_dates.append(ed.strftime('%Y-%m-%d') if ed < today_ts else '')
        else:
            end_dates.append('')
    revenues = np.random.choice([300, 500, 980], n_sample, p=[0.5, 0.35, 0.15])
    sample = pd.DataFrame({
        'customer_id':     [f'C{i:05d}' for i in range(1, n_sample+1)],
        'start_date':      [d.strftime('%Y-%m-%d') for d in start_dates],
        'end_date':        end_dates,
        'revenue_monthly': revenues,
    })
    st.download_button(
        "サンプルCSV（1万件）をダウンロード",
        sample.to_csv(index=False).encode('utf-8-sig'),
        "sample_customers_10000.csv", "text/csv"
    )
    st.caption("列名 `revenue_monthly` は月次売上（円）のサンプルです")

    uploaded = st.file_uploader("CSVをアップロード", type=['csv'])

    st.markdown("---")
    st.markdown("### ⚙️ 設定")
    st.markdown("**CSVの売上列の単位を選択してください**")
    st.caption("ツール内部で自動的に日次ARPUに変換して計算します")
    revenue_unit_raw = st.selectbox(
        "売上単位",
        ['日次（1日あたりの売上）', '月次（1ヶ月あたりの売上）', '年次（1年あたりの売上）'],
        index=1,
        label_visibility='collapsed'
    )
    revenue_unit = revenue_unit_raw.split('（')[0]
    horizon_days = st.number_input("観測期間上限（日）", 30, 3650, 730, help="730日=2年、365日=1年が目安")

    st.markdown("---")
    st.markdown("### 📊 粗利率（GPM）")
    gpm = st.slider("Gross Profit Margin (%)", 0, 100, 54, 1) / 100
    st.caption(f"LTVは売上ベースではなく粗利ベース（×{gpm:.0%}）で算出されます")

    st.markdown("---")
    st.markdown("### 💰 CAC上限の算出")
    cac_mode = st.radio("算出方法", ['LTV ÷ N', 'LTV : CAC = N : 1', '回収期間（月）'])
    if cac_mode == 'LTV ÷ N':
        cac_n = st.slider("N（分母）", 1.0, 10.0, 3.0, 0.5)
        cac_label = f"LTV ÷ {cac_n}"
    elif cac_mode == 'LTV : CAC = N : 1':
        cac_n = st.slider("N", 1.0, 10.0, 3.0, 0.5)
        cac_label = f"LTV:CAC = {cac_n}:1"
    else:
        cac_n = st.slider("回収期間（月）", 1, 36, 12)
        cac_label = f"{cac_n}ヶ月回収"

    st.markdown("---")
    st.markdown("### 🏢 レポート情報")
    client_name  = st.text_input("クライアント名", "")
    analyst_name = st.text_input("分析者名", "")

# ══════════════════════════════════════════════════════════════
# Header
# ══════════════════════════════════════════════════════════════

st.markdown("""
<div style='padding: 8px 0 24px 0;'>
  <span style='font-family:Syne,sans-serif; font-size:2rem; color:#e8e4dc;'>LTV Analyzer</span>
  <span style='font-size:0.8rem; color:#444; margin-left:12px;'>Kaplan–Meier × Weibull Model</span>
</div>
""", unsafe_allow_html=True)

# ══════════════════════════════════════════════════════════════
# No file → instructions
# ══════════════════════════════════════════════════════════════

if uploaded is None:
    st.info("← サイドバーからCSVをアップロードしてください。サンプルCSVでまずお試しいただけます。")

    st.markdown("<div class='section-title'>入力CSVの形式</div>", unsafe_allow_html=True)
    st.markdown("""
| 列名 | 内容 | 形式 | 例 |
|------|------|------|----|
| `customer_id` | 顧客ID | 任意の文字列 | C0001 |
| `start_date` | 契約開始日 | YYYY-MM-DD | 2023-10-01 |
| `end_date` | 解約日（継続中は**空欄**） | YYYY-MM-DD | 2024-03-15 |
| `revenue` | 売上（サイドバーで単位を選択） | 数値 | 300 |

> 列名は完全一致でなくてもOKです。`start`・`end`・`id`・`revenue`を含む列名は自動認識します。
    """)

    st.markdown("<div class='section-title'>分析の流れ</div>", unsafe_allow_html=True)
    cols = st.columns(4)
    steps = [
        ("① KM法", "実測データから生存曲線を作成"),
        ("② Weibull", "連続曲線にフィッティング"),
        ("③ LTV∞", "生存積分 × ARPU で算出"),
        ("④ CAC上限", "LTV比率で逆算"),
    ]
    for col, (title, desc) in zip(cols, steps):
        with col:
            st.markdown(f"<div class='metric-card'><div class='metric-value' style='font-size:1.1rem'>{title}</div><div class='metric-label' style='font-size:0.75rem; color:#666; margin-top:8px;'>{desc}</div></div>", unsafe_allow_html=True)
    st.stop()

# ══════════════════════════════════════════════════════════════
# Load & validate data
# ══════════════════════════════════════════════════════════════

try:
    df_raw = pd.read_csv(uploaded)
    df_raw.columns = df_raw.columns.str.strip().str.lower()

    col_map = {}
    for c in df_raw.columns:
        if any(x in c for x in ['id','customer','顧客']):      col_map.setdefault('customer_id', c)
        if any(x in c for x in ['start','開始','契約']):       col_map.setdefault('start_date', c)
        if any(x in c for x in ['end','解約','終了','cancel']): col_map.setdefault('end_date', c)
        if any(x in c for x in ['revenue','売上','arpu','rev','price','amount']):
            col_map.setdefault('revenue', c)

    missing = [k for k in ['start_date','end_date','revenue'] if k not in col_map]
    if missing:
        st.error(f"❌ 列が見つかりません: {missing}\n\n列名に `start`・`end`・`revenue` を含む列が必要です。サイドバーからサンプルCSVをダウンロードして形式を確認してください。")
        st.stop()

    df = df_raw.rename(columns={v: k for k, v in col_map.items()})
    df['start_date'] = pd.to_datetime(df['start_date'], errors='coerce')
    df['end_date']   = pd.to_datetime(df['end_date'], errors='coerce')

    bad_dates = df['start_date'].isna().sum()
    if bad_dates > 0:
        st.warning(f"⚠️ {bad_dates}行で `start_date` が読み取れませんでした。該当行は除外します。")
    df = df.dropna(subset=['start_date'])

    today = pd.Timestamp.today()
    n_input = len(df)
    df['duration'] = (df['end_date'].fillna(today) - df['start_date']).dt.days
    df['event']    = df['end_date'].notna().astype(int)

    # duration=0 を1日に自動補正（入力件数＝分析件数を保証）
    n_corrected = (df['duration'] == 0).sum()
    df['duration'] = df['duration'].clip(lower=1)

    # duration<0（開始日が未来）は除外
    n_excluded = (df['duration'] < 0).sum()
    df = df[df['duration'] > 0]

    if n_corrected > 0:
        st.info(f"ℹ️ {n_corrected}件のデータで契約開始日と解約日が同日でした。自動的に1日に補正しました。")
    if n_excluded > 0:
        st.warning(f"⚠️ {n_excluded}件のデータで契約開始日が未来の日付でした。該当行を除外しました。")
    if n_corrected == 0 and n_excluded == 0:
        st.success(f"✅ 全{n_input:,}件のデータを正常に読み込みました。")

    if len(df) < 10:
        st.error("❌ 有効なデータが10件未満です。分析には最低10件の顧客データが必要です。")
        st.stop()

    unit_div = {'日次': 1, '月次': 30.44, '年次': 365}
    df['arpu_daily'] = pd.to_numeric(df['revenue'], errors='coerce') / unit_div[revenue_unit]
    df['gp_daily']   = df['arpu_daily'] * gpm
    df = df.dropna(subset=['arpu_daily'])
    arpu_daily = df['arpu_daily'].mean()
    gp_daily   = df['gp_daily'].mean()

except Exception as e:
    st.error(f"❌ データ読み込みエラー: {e}\n\nCSVの形式を確認してください。サンプルCSVをダウンロードして参照してください。")
    st.stop()

# ══════════════════════════════════════════════════════════════
# Run analysis
# ══════════════════════════════════════════════════════════════

km_df = compute_km(df)
k, lam, r2, fit_df = fit_weibull(km_df)

if k is None:
    st.error("❌ Weibullフィッティングに失敗しました。解約済み顧客が少なすぎる可能性があります（最低10件の解約データが必要）。")
    st.stop()

ltv_rev, surv_int = ltv_inf(k, lam, arpu_daily)   # 売上ベース
ltv_val, _        = ltv_inf(k, lam, gp_daily)      # 粗利ベース（CACに使う）
cac_upper = ltv_val / cac_n

# ══════════════════════════════════════════════════════════════
# Metrics
# ══════════════════════════════════════════════════════════════

st.markdown("<div class='section-title'>分析結果サマリー</div>", unsafe_allow_html=True)
m1, m2, m3, m4, m5 = st.columns(5)

metrics = [
    (f"¥{ltv_rev:,.0f}", "LTV∞（売上ベース）"),
    (f"¥{cac_upper:,.0f}", f"CAC上限 {cac_label}"),
    (f"{k:.3f}", "Weibull k（形状） k<1:初期離脱多い / k>1:時間とともに離脱増"),
    (f"{lam:.1f}日", "Weibull λ（尺度） 顧客の典型的な生存日数の目安"),
    (f"{r2:.3f}", "R² 1.0に近いほど精度高い（0.9以上が理想）"),
]
for col, (val, label) in zip([m1,m2,m3,m4,m5], metrics):
    with col:
        st.markdown(f"<div class='metric-card'><div class='metric-value'>{val}</div><div class='metric-label'>{label}</div></div>", unsafe_allow_html=True)

# R² warning
if r2 < 0.85:
    st.warning(f"⚠️ R²={r2:.3f} — フィット精度がやや低めです。データ点数を増やすか、観測期間を見直してください。")

# ══════════════════════════════════════════════════════════════
# Charts
# ══════════════════════════════════════════════════════════════

st.markdown("<div class='section-title'>分析グラフ</div>", unsafe_allow_html=True)
c1, c2 = st.columns(2)

t_smooth = np.linspace(1, km_df['t'].max() * 1.3, 600)
S_wei    = weibull_s(t_smooth, k, lam)

with c1:
    fig, ax = plt.subplots(figsize=(6, 3.8))
    ax.step(km_df['t'], km_df['S'], where='post', color=ACCENT, lw=1.8, label='KM生存曲線（実測）')
    ax.plot(t_smooth, S_wei, color=ACCENT2, lw=1.5, ls='--', label='Weibullフィット')
    ax.fill_between(t_smooth, S_wei, alpha=0.06, color=ACCENT2)
    ax.set(xlabel='生存日数（日）', ylabel='生存率 S(t)', ylim=(0,1.05))
    ax.legend(fontsize=8, framealpha=0.15)
    ax.grid(True, alpha=0.25)
    ax.set_title('生存曲線', color='#ccc', fontsize=10, pad=8)
    fig.tight_layout()
    st.pyplot(fig)
    plt.close()

with c2:
    fig, ax = plt.subplots(figsize=(6, 3.8))
    if fit_df is not None:
        ax.scatter(fit_df['ln_t'], fit_df['ln_neg_ln_S'], color=ACCENT, s=18, alpha=0.7, label='実測点')
        x_r = np.linspace(fit_df['ln_t'].min(), fit_df['ln_t'].max(), 200)
        b   = np.log(-np.log(weibull_s(1, k, lam) + 1e-15))
        ax.plot(x_r, k*x_r + b, color=ACCENT2, lw=1.5, ls='--', label=f'回帰直線 (R²={r2:.3f})')
        ax.annotate(f'y = {k:.4f}x + {b:.4f}',
                    xy=(0.05,0.93), xycoords='axes fraction', color='#777', fontsize=8)
    ax.set(xlabel='ln(t)', ylabel='ln(−ln(S(t)))')
    ax.legend(fontsize=8, framealpha=0.15)
    ax.grid(True, alpha=0.25)
    ax.set_title('Weibull直線化プロット', color='#ccc', fontsize=10, pad=8)
    fig.tight_layout()
    st.pyplot(fig)
    plt.close()

# Save chart images for export
fig1, ax1 = plt.subplots(figsize=(7, 4))
ax1.step(km_df['t'], km_df['S'], where='post', color=ACCENT, lw=2, label='KM生存曲線')
ax1.plot(t_smooth, S_wei, color=ACCENT2, lw=1.8, ls='--', label='Weibullフィット')
ax1.fill_between(t_smooth, S_wei, alpha=0.07, color=ACCENT2)
ax1.set(xlabel='生存日数（日）', ylabel='生存率 S(t)', ylim=(0,1.05))
ax1.legend(fontsize=9, framealpha=0.15)
ax1.grid(True, alpha=0.25)
ax1.set_title('生存曲線（KM × Weibull）', color='#ccc', fontsize=11, pad=10)
fig1.tight_layout()
buf1 = io.BytesIO(); fig1.savefig(buf1, format='png', dpi=150, bbox_inches='tight'); buf1.seek(0)
plt.close()

fig2, ax2 = plt.subplots(figsize=(7, 4))
if fit_df is not None:
    ax2.scatter(fit_df['ln_t'], fit_df['ln_neg_ln_S'], color=ACCENT, s=22, alpha=0.75)
    x_r = np.linspace(fit_df['ln_t'].min(), fit_df['ln_t'].max(), 200)
    b   = np.log(-np.log(weibull_s(1, k, lam) + 1e-15))
    ax2.plot(x_r, k*x_r+b, color=ACCENT2, lw=1.8, ls='--', label=f'R²={r2:.3f}')
    ax2.annotate(f'y = {k:.4f}x + {b:.4f}', xy=(0.05,0.93), xycoords='axes fraction', color='#777', fontsize=9)
ax2.set(xlabel='ln(t)', ylabel='ln(−ln(S(t)))')
ax2.legend(fontsize=9, framealpha=0.15)
ax2.grid(True, alpha=0.25)
ax2.set_title('Weibull直線化プロット', color='#ccc', fontsize=11, pad=10)
fig2.tight_layout()
buf2 = io.BytesIO(); fig2.savefig(buf2, format='png', dpi=150, bbox_inches='tight'); buf2.seek(0)
plt.close()

# ══════════════════════════════════════════════════════════════
# Horizon table
# ══════════════════════════════════════════════════════════════

st.markdown("<div class='section-title'>暫定LTV（観測期間別）</div>", unsafe_allow_html=True)

horizons = sorted(set([30, 90, 180, 365, 730, int(horizon_days)]))
rows = []
for h in horizons:
    lh = ltv_horizon(k, lam, arpu_daily, h)
    rows.append({
        'ホライズン': f'{h}日' if h < 365 else (f'{h//365}年' if h%365==0 else f'{h}日'),
        '暫定LTV': f'¥{lh:,.0f}',
        'LTV∞比': f'{lh/ltv_val*100:.1f}%',
        f'CAC上限 ({cac_label})': f'¥{lh/cac_n:,.0f}',
    })
st.dataframe(pd.DataFrame(rows), hide_index=True, use_container_width=True)

# ══════════════════════════════════════════════════════════════
# AI Prompt Generator
# ══════════════════════════════════════════════════════════════

st.markdown("<div class='section-title'>🤖 AIに質問するプロンプト</div>", unsafe_allow_html=True)
st.markdown("<div class='help-box'>この結果の読み方や戦略への活用方法がわからない場合は、以下のプロンプトをClaude・ChatGPT・Geminiにコピペしてください。</div>", unsafe_allow_html=True)

tab1, tab2, tab3 = st.tabs(["📊 結果の読み方", "📈 マーケ戦略への活用", "⚠️ 精度の検証"])

prompt_base = f"""私はLTV分析ツールを使い、以下の結果を得ました。

【分析結果】
・顧客数: {len(df):,}件（うち解約済み: {df['event'].sum():,}件）
・平均日次ARPU（売上ベース）: ¥{arpu_daily:,.2f}
・粗利率（GPM）: {gpm:.1%}
・平均日次GP（粗利ベース）: ¥{gp_daily:,.2f}
・LTV∞（売上ベース）: ¥{ltv_rev:,.0f}
・LTV∞（粗利ベース）: ¥{ltv_val:,.0f}
・CAC上限 ({cac_label}): ¥{cac_upper:,.0f}
・Weibull 形状パラメータ k: {k:.4f}
・Weibull 尺度パラメータ λ: {lam:.1f}日
・R²（フィット精度）: {r2:.4f}
・分析手法: Kaplan-Meier法 + Weibullモデルによる生存分析"""

with tab1:
    p1 = prompt_base + """

【質問】
1. Weibullのkとλの値は何を意味していますか？このビジネスの顧客離脱パターンはどう解釈すればよいですか？
2. LTV∞の値は適切な水準ですか？
3. R²の値からフィット精度はどう評価できますか？
4. この結果で特に注意すべき点があれば教えてください。"""
    st.markdown(f"<div class='prompt-box'>{p1}</div>", unsafe_allow_html=True)
    st.button("📋 コピー（結果の読み方）", on_click=lambda: st.write(""), key="copy1",
              help="上のテキストを選択してCtrl+CでコピーできますClaude/ChatGPT/Geminiに貼り付けてください")

with tab2:
    p2 = prompt_base + f"""
・観測期間: {horizon_days}日

【質問】
1. このLTV∞とCAC上限をもとに、広告予算の上限をどう設定すべきですか？
2. 顧客獲得チャネル別にROIを評価するには何が必要ですか？
3. LTVを高めるために優先すべき施策は何ですか？
4. このビジネスに最適なLTV:CAC比率の目安を教えてください。"""
    st.markdown(f"<div class='prompt-box'>{p2}</div>", unsafe_allow_html=True)

with tab3:
    p3 = prompt_base + """

【質問】
1. このデータ件数と解約件数でWeibullフィッティングの信頼性はどう評価できますか？
2. R²の値は十分ですか？改善するにはどうすればよいですか？
3. Weibullモデルの仮定が成立していない可能性はありますか？どうチェックすればよいですか？
4. セグメント分割（プラン別・属性別）をするメリットとデメリットを教えてください。"""
    st.markdown(f"<div class='prompt-box'>{p3}</div>", unsafe_allow_html=True)

# ══════════════════════════════════════════════════════════════
# Export buttons
# ══════════════════════════════════════════════════════════════

st.markdown("<div class='section-title'>📤 エクスポート</div>", unsafe_allow_html=True)

exp1, exp2, exp3 = st.columns(3)

# ── Excel export ─────────────────────────────────────────────
with exp1:
    try:
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
        from openpyxl.chart import LineChart, Reference, Series

        wb = openpyxl.Workbook()

        # Summary sheet
        ws = wb.active
        ws.title = 'Summary'
        hdr_fill = PatternFill('solid', start_color='1a1a1a', end_color='1a1a1a')
        hdr_font = Font(name='Calibri', bold=True, color='C8B89A', size=11)
        val_font = Font(name='Calibri', size=11, color='E8E4DC')

        ws['A1'] = 'LTV分析 サマリー'
        ws['A1'].font = Font(name='Calibri', bold=True, size=14, color='C8B89A')
        if client_name:
            ws['A2'] = f'クライアント: {client_name}'
        if analyst_name:
            ws['A3'] = f'分析者: {analyst_name}'

        summary_data = [
            ('', ''),
            ('【分析結果】', ''),
            ('顧客数（総数）', len(df)),
            ('うち解約済み', int(df['event'].sum())),
            ('平均日次ARPU（売上ベース・¥）', round(arpu_daily, 2)),
            ('粗利率（GPM）', f'{gpm:.1%}'),
            ('平均日次GP（粗利ベース・¥）', round(gp_daily, 2)),
            ('LTV∞（売上ベース・¥）', round(ltv_rev, 0)),
            ('LTV∞（粗利ベース・¥）', round(ltv_val, 0)),
            (f'CAC上限（{cac_label}）（¥）', round(cac_upper, 0)),
            ('', ''),
            ('【Weibullパラメータ】', ''),
            ('k（形状パラメータ）', round(k, 4)),
            ('λ（尺度パラメータ・日）', round(lam, 2)),
            ('R²', round(r2, 4)),
        ]
        for i, (label, val) in enumerate(summary_data, start=5):
            ws.cell(i, 1, label).font = Font(name='Calibri', bold=('【' in str(label)), color='888888', size=10)
            ws.cell(i, 2, val).font   = Font(name='Calibri', size=10, color='E8E4DC')
        ws.column_dimensions['A'].width = 32
        ws.column_dimensions['B'].width = 20

        # KM sheet
        ws2 = wb.create_sheet('KM_生存曲線')
        ws2.append(['t（日）', 'S(t)_KM実測', 'S(t)_Weibullフィット'])
        for _, row in km_df.iterrows():
            t = row['t']
            ws2.append([int(t), round(row['S'], 6), round(float(weibull_s(t, k, lam)), 6)])

        # Horizon sheet
        ws3 = wb.create_sheet('暫定LTV')
        ws3.append(['ホライズン（日）', '暫定LTV（¥）', 'LTV∞比（%）', f'CAC上限（¥）'])
        for h in horizons:
            lh = ltv_horizon(k, lam, arpu_daily, h)
            ws3.append([h, round(lh, 0), round(lh/ltv_val*100, 1), round(lh/cac_n, 0)])

        xl_buf = io.BytesIO()
        wb.save(xl_buf)
        xl_buf.seek(0)
        st.download_button("📊 Excelダウンロード", xl_buf,
                           file_name=f"LTV分析_{client_name or 'report'}.xlsx",
                           mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
    except Exception as e:
        st.caption(f"Excel出力エラー: {e}")

# ── PowerPoint export ─────────────────────────────────────────
with exp2:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        BG    = RGBColor(0x0d, 0x0d, 0x0d)
        GOLD  = RGBColor(0xc8, 0xb8, 0x9a)
        TEAL  = RGBColor(0x7a, 0x9e, 0x9f)
        WHITE = RGBColor(0xe8, 0xe4, 0xdc)
        GRAY  = RGBColor(0x44, 0x44, 0x44)

        def add_bg(slide, prs):
            bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
            bg.fill.solid(); bg.fill.fore_color.rgb = BG
            bg.line.fill.background(); bg.zorder = 0

        def txbox(slide, text, l, t, w, h, size=12, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
            tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
            tf = tb.text_frame; tf.word_wrap = True
            p  = tf.paragraphs[0]; p.alignment = align
            run = p.add_run(); run.text = text
            run.font.size  = Pt(size)
            run.font.bold  = bold
            run.font.color.rgb = color
            return tb

        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]

        # ── Slide 1: Title ──
        s1 = prs.slides.add_slide(blank)
        add_bg(s1, prs)
        # Accent line
        line = s1.shapes.add_shape(1, Inches(0.8), Inches(3.2), Inches(0.04), Inches(1.4))
        line.fill.solid(); line.fill.fore_color.rgb = GOLD; line.line.fill.background()
        txbox(s1, 'LTV Analysis Report', 1.0, 2.8, 8, 0.8, size=36, bold=True, color=WHITE)
        txbox(s1, 'Kaplan–Meier × Weibull Model', 1.0, 3.7, 8, 0.5, size=14, color=GOLD)
        if client_name:
            txbox(s1, client_name, 1.0, 4.4, 8, 0.4, size=13, color=RGBColor(0x88,0x88,0x88))
        from datetime import date
        txbox(s1, date.today().strftime('%Y年%m月%d日'), 1.0, 5.0, 6, 0.4, size=11, color=GRAY)
        if analyst_name:
            txbox(s1, analyst_name, 1.0, 5.4, 6, 0.4, size=11, color=GRAY)

        # ── Slide 2: Metrics ──
        s2 = prs.slides.add_slide(blank)
        add_bg(s2, prs)
        txbox(s2, '分析結果サマリー', 0.5, 0.3, 10, 0.6, size=22, bold=True, color=WHITE)

        cards = [
            ('LTV∞', f'¥{ltv_val:,.0f}', 0.5),
            (f'CAC上限\n({cac_label})', f'¥{cac_upper:,.0f}', 3.6),
            ('Weibull k', f'{k:.3f}', 6.7),
            ('R²', f'{r2:.3f}', 9.8),
        ]
        for label, val, x in cards:
            card = s2.shapes.add_shape(1, Inches(x), Inches(1.3), Inches(2.9), Inches(1.8))
            card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0x16,0x16,0x16)
            card.line.color.rgb = GRAY
            txbox(s2, val, x+0.15, 1.45, 2.6, 0.8, size=24, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
            txbox(s2, label, x+0.15, 2.2, 2.6, 0.5, size=10, color=GRAY, align=PP_ALIGN.CENTER)

        # Data stats
        stats_text = (
            f"顧客数: {len(df):,}件  ／  解約済み: {df['event'].sum():,}件  ／  "
            f"継続中: {(df['event']==0).sum():,}件  ／  "
            f"平均日次ARPU: ¥{arpu_daily:,.2f}  ／  λ: {lam:.1f}日"
        )
        txbox(s2, stats_text, 0.5, 3.3, 12.3, 0.5, size=10, color=GRAY)

        # Horizon table header
        txbox(s2, '観測期間別 暫定LTV', 0.5, 3.9, 12, 0.4, size=13, bold=True, color=WHITE)
        cols_h = ['ホライズン', '暫定LTV', 'LTV∞比', 'CAC上限']
        col_x  = [0.5, 3.3, 6.5, 9.3]
        for cx, ch in zip(col_x, cols_h):
            txbox(s2, ch, cx, 4.4, 2.6, 0.35, size=9, bold=True, color=GOLD)
        row_y = 4.8
        for h in horizons[:5]:
            lh = ltv_horizon(k, lam, arpu_daily, h)
            row_vals = [
                f'{h}日' if h<365 else f'{h//365}年',
                f'¥{lh:,.0f}', f'{lh/ltv_val*100:.1f}%', f'¥{lh/cac_n:,.0f}'
            ]
            for cx, rv in zip(col_x, row_vals):
                txbox(s2, rv, cx, row_y, 2.6, 0.32, size=9, color=WHITE)
            row_y += 0.34

        # ── Slide 3: Charts ──
        s3 = prs.slides.add_slide(blank)
        add_bg(s3, prs)
        txbox(s3, '生存曲線 / Weibull直線化プロット', 0.5, 0.3, 12, 0.6, size=22, bold=True, color=WHITE)
        buf1.seek(0); s3.shapes.add_picture(buf1, Inches(0.4), Inches(1.1), Inches(6.1), Inches(3.8))
        buf2.seek(0); s3.shapes.add_picture(buf2, Inches(6.8), Inches(1.1), Inches(6.1), Inches(3.8))

        note_text = f"k={k:.3f} (形状) ／ λ={lam:.1f}日 (尺度) ／ R²={r2:.3f}  ※ k<1: 初期離脱が多い、k>1: 時間とともに離脱が増加"
        txbox(s3, note_text, 0.4, 5.1, 12.5, 0.5, size=9, color=GRAY)

        # ── Slide 4: AI Prompt ──
        s4 = prs.slides.add_slide(blank)
        add_bg(s4, prs)
        txbox(s4, 'AIへの質問プロンプト', 0.5, 0.3, 12, 0.6, size=22, bold=True, color=WHITE)
        txbox(s4, 'この結果についてClaude / ChatGPT / Geminiに質問する際のプロンプトです。コピペしてご活用ください。', 0.5, 1.0, 12, 0.4, size=10, color=GRAY)

        prompt_short = (
            f"LTV分析の結果です。\n"
            f"顧客数: {len(df):,}件、LTV∞: ¥{ltv_val:,.0f}、CAC上限({cac_label}): ¥{cac_upper:,.0f}\n"
            f"Weibull k={k:.4f}、λ={lam:.1f}日、R²={r2:.4f}\n\n"
            f"1. kとλの値はどう解釈すればよいですか？\n"
            f"2. このLTVとCACをもとに広告予算をどう設定すべきですか？\n"
            f"3. LTVを高めるために優先すべき施策は何ですか？"
        )
        pbox = s4.shapes.add_shape(1, Inches(0.5), Inches(1.6), Inches(12.3), Inches(4.5))
        pbox.fill.solid(); pbox.fill.fore_color.rgb = RGBColor(0x14,0x14,0x14)
        pbox.line.color.rgb = GOLD
        txbox(s4, prompt_short, 0.7, 1.75, 11.9, 4.2, size=10, color=WHITE)

        pptx_buf = io.BytesIO()
        prs.save(pptx_buf)
        pptx_buf.seek(0)
        st.download_button("📑 PowerPointダウンロード", pptx_buf,
                           file_name=f"LTV分析_{client_name or 'report'}.pptx",
                           mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
    except ImportError:
        st.caption("PowerPoint出力には `pip install python-pptx` が必要です")
    except Exception as e:
        st.caption(f"PowerPoint出力エラー: {e}")

# ── PDF export ────────────────────────────────────────────────
with exp3:
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib import colors
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, Table, TableStyle
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import cm
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.cidfonts import UnicodeCIDFont
        pdfmetrics.registerFont(UnicodeCIDFont('HeiseiMin-W3'))

        pdf_buf = io.BytesIO()
        doc = SimpleDocTemplate(pdf_buf, pagesize=A4,
                                leftMargin=2*cm, rightMargin=2*cm,
                                topMargin=2*cm, bottomMargin=2*cm)
        styles = getSampleStyleSheet()
        title_style = ParagraphStyle('T', fontName='HeiseiMin-W3', fontSize=18, spaceAfter=6,
                                     textColor=colors.HexColor('#c8b89a'))
        h2_style    = ParagraphStyle('H2', fontName='HeiseiMin-W3', fontSize=13, spaceAfter=4,
                                     textColor=colors.HexColor('#e8e4dc'), spaceBefore=14)
        body_style  = ParagraphStyle('B', fontName='HeiseiMin-W3', fontSize=9,
                                     textColor=colors.HexColor('#888888'), spaceAfter=3)

        story = []
        story.append(Paragraph('LTV Analysis Report', title_style))
        if client_name: story.append(Paragraph(f'クライアント: {client_name}', body_style))
        story.append(Spacer(1, 0.3*cm))

        story.append(Paragraph('分析結果サマリー', h2_style))
        tdata = [
            ['指標', '値'],
            ['LTV∞', f'¥{ltv_val:,.0f}'],
            [f'CAC上限 ({cac_label})', f'¥{cac_upper:,.0f}'],
            ['Weibull k（形状）', f'{k:.4f}'],
            ['Weibull λ（尺度）', f'{lam:.1f}日'],
            ['R²', f'{r2:.4f}'],
            ['顧客数', f'{len(df):,}件'],
            ['解約済み', f'{int(df["event"].sum()):,}件'],
            ['平均日次ARPU', f'¥{arpu_daily:,.2f}'],
        ]
        t = Table(tdata, colWidths=[9*cm, 6*cm])
        t.setStyle(TableStyle([
            ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#1a1a1a')),
            ('TEXTCOLOR',  (0,0), (-1,0), colors.HexColor('#c8b89a')),
            ('TEXTCOLOR',  (0,1), (-1,-1), colors.HexColor('#cccccc')),
            ('FONTNAME',   (0,0), (-1,-1), 'HeiseiMin-W3'),
            ('FONTSIZE',   (0,0), (-1,-1), 9),
            ('ROWBACKGROUNDS', (0,1), (-1,-1), [colors.HexColor('#111111'), colors.HexColor('#161616')]),
            ('GRID', (0,0), (-1,-1), 0.3, colors.HexColor('#333333')),
            ('LEFTPADDING', (0,0), (-1,-1), 8),
        ]))
        story.append(t)

        story.append(Paragraph('生存曲線', h2_style))
        buf1.seek(0)
        story.append(Image(buf1, width=15*cm, height=9*cm))

        story.append(Paragraph('Weibull直線化プロット', h2_style))
        buf2.seek(0)
        story.append(Image(buf2, width=15*cm, height=9*cm))

        doc.build(story)
        pdf_buf.seek(0)
        st.download_button("📄 PDFダウンロード", pdf_buf,
                           file_name=f"LTV分析_{client_name or 'report'}.pdf",
                           mime="application/pdf")
    except ImportError:
        st.caption("PDF出力には `pip install reportlab` が必要です")
    except Exception as e:
        st.caption(f"PDF出力エラー: {e}")

# ══════════════════════════════════════════════════════════════
# Data preview
# ══════════════════════════════════════════════════════════════

with st.expander("📋 読み込んだデータを確認"):
    st.write(f"有効データ: {len(df):,}件 ／ 解約: {df['event'].sum():,}件 ／ 継続中: {(df['event']==0).sum():,}件 ／ 平均日次ARPU: ¥{arpu_daily:,.2f}")
    st.dataframe(
        df[['customer_id','start_date','end_date','duration','event','arpu_daily']].head(30),
        hide_index=True
    )

st.markdown("---")
st.markdown("<p style='color:#333; font-size:0.72rem; text-align:center;'>LTV Analyzer — KM × Weibull Model — Built for marketing analytics professionals</p>", unsafe_allow_html=True)
